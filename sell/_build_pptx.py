# -*- coding: utf-8 -*-
"""Build 16:9 PPTX decks from sell/<product>/deck.json."""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.util import Emu, Inches, Pt
from lxml import etree

INK = RGBColor(0x1A, 0x18, 0x14)
PAPER = RGBColor(0xF4, 0xF0, 0xE6)
RUST = RGBColor(0xC4, 0x5C, 0x26)
CREAM = RGBColor(0xE8, 0xE0, 0xD0)
MUTED = RGBColor(0x6A, 0x62, 0x56)
WHITE = RGBColor(0xF7, 0xF3, 0xEC)

W = Inches(13.333)
H = Inches(7.5)


def font_for(lang: str) -> str:
    if lang == "zh":
        return "Microsoft YaHei"
    return "Calibri"


def set_run(run, text, size, bold, color, lang):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    face = font_for(lang)
    run.font.name = face
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(ns + "ea")
    if ea is None:
        ea = etree.SubElement(rPr, ns + "ea")
    ea.set("typeface", face)


def add_textbox(slide, l, t, w, h, text, size, color, lang, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size, bold, color, lang)
    return box


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_notes(slide, note: str):
    if not note:
        return
    ns = slide.notes_slide.notes_text_frame
    ns.text = note


def paint_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill_shape(bg, color)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def rust_bar(slide, dark=False):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
    fill_shape(bar, RUST if not dark else RGBColor(0xE0, 0x7A, 0x3C))


def cover_slide(prs, block, lang, product_footer):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, INK)
    rust_bar(s, dark=True)
    k = block.get("kicker") or ""
    add_textbox(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.45), k.upper(), 13, RUST, lang, bold=True)
    add_textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.2), block.get("title") or "", 40, WHITE, lang, bold=True)
    lines = block.get("lines") or []
    y = 4.4
    for line in lines[:3]:
        add_textbox(s, Inches(0.7), Inches(y), Inches(11.5), Inches(0.4), line, 18, CREAM, lang)
        y += 0.42
    add_textbox(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.35), product_footer, 12, MUTED, lang)
    add_notes(s, block.get("note") or "")


def content_slide(prs, block, lang, idx, n, footer):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, PAPER)
    rust_bar(s)
    k = block.get("kicker") or ""
    if k:
        add_textbox(s, Inches(0.7), Inches(0.35), Inches(11), Inches(0.35), k.upper(), 12, RUST, lang, bold=True)
    add_textbox(s, Inches(0.7), Inches(0.7), Inches(12), Inches(1.1), block.get("title") or "", 28, INK, lang, bold=True)
    lines = block.get("lines") or []
    y = 2.05
    for line in lines:
        mark = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y + 0.12), Inches(0.12), Inches(0.12))
        fill_shape(mark, RUST)
        add_textbox(s, Inches(1.05), Inches(y), Inches(11.4), Inches(0.85), line, 18, INK, lang)
        y += 0.78
        if y > 6.3:
            break
    add_textbox(s, Inches(0.7), Inches(7.05), Inches(10), Inches(0.3), footer, 11, MUTED, lang)
    add_textbox(s, Inches(11.6), Inches(7.05), Inches(1.3), Inches(0.3), "%d / %d" % (idx, n), 11, MUTED, lang, align=PP_ALIGN.RIGHT)
    add_notes(s, block.get("note") or "")


def repro_slide(prs, block, lang, idx, n, footer):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, PAPER)
    rust_bar(s)
    k = block.get("kicker") or ""
    if k:
        add_textbox(s, Inches(0.7), Inches(0.35), Inches(11), Inches(0.35), k.upper(), 12, RUST, lang, bold=True)
    add_textbox(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.9), block.get("title") or "", 28, INK, lang, bold=True)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.7))
    fill_shape(panel, INK)
    tf = s.shapes.add_textbox(Inches(1.0), Inches(2.05), Inches(11.3), Inches(4.3)).text_frame
    tf.word_wrap = True
    first = True
    for line in block.get("lines") or []:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        run = p.add_run()
        set_run(run, line, 16, False, CREAM, lang)
        run.font.name = "Consolas" if lang != "zh" else font_for(lang)
    add_textbox(s, Inches(0.7), Inches(7.05), Inches(10), Inches(0.3), footer, 11, MUTED, lang)
    add_textbox(s, Inches(11.6), Inches(7.05), Inches(1.3), Inches(0.3), "%d / %d" % (idx, n), 11, MUTED, lang, align=PP_ALIGN.RIGHT)
    add_notes(s, block.get("note") or "")


def build_one(deck_path: Path, lang: str, out: Path):
    data = json.loads(deck_path.read_text(encoding="utf-8"))
    slides = data["slides"]
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    footer = {
        "ru": "Kenga | not a press kit | numbers from the repo",
        "en": "Kenga | not a press kit | numbers from the repo only",
        "zh": "Kenga | not a press kit | numbers from the repo only",
    }[lang]
    n = len(slides)
    for i, sl in enumerate(slides, 1):
        block = sl[lang]
        typ = sl.get("type") or ""
        if typ == "cover" or i == 1:
            cover_slide(prs, block, lang, footer)
        elif typ == "repro":
            repro_slide(prs, block, lang, i, n, footer)
        else:
            content_slide(prs, block, lang, i, n, footer)
    prs.save(str(out))
    print("wrote", out)


def main():
    root = Path(r"d:\kenga-lang\sell")
    products = [
        "kenga-language",
        "kenga-factory",
        "kenga-expert-lm",
        "kenga-baremetal",
        "kenga-passport",
    ]
    langs = ["ru", "en", "zh"]
    for p in products:
        deck = root / p / "deck.json"
        if not deck.exists():
            print("missing", deck)
            continue
        for lang in langs:
            build_one(deck, lang, root / p / ("pitch.%s.pptx" % lang))


if __name__ == "__main__":
    main()
